"""Generate an editable PowerPoint .pptx version of the SGP-Pharma user manual."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

SCREENSHOTS = Path("/app/docs/screenshots")
OUT = Path("/app/docs/SGP-Pharma_Manuel_Utilisateur.pptx")

PRIMARY = RGBColor(0x16, 0x65, 0x34)
PRIMARY_DARK = RGBColor(0x14, 0x53, 0x2D)
TEXT = RGBColor(0x11, 0x18, 0x27)
MUTED = RGBColor(0x4B, 0x55, 0x63)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xF5, 0x9E, 0x0B)
LIGHT_BG = RGBColor(0xF9, 0xFA, 0xFB)


SLIDES = [
    # (badge, title, subtitle, bullets, screenshot)
    ("MODULE 1 — AUTHENTIFICATION",
     "Connexion & Gestion des rôles",
     "Authentification sécurisée par JWT — 5 niveaux d'accès distincts.",
     [
         ("Comptes de démonstration", True),
         ("Super Admin (OPTINET) — optinet@sgp-pharma.tg", False),
         ("Admin Pharmacie — admin@sgp-pharma.tg", False),
         ("Pharmacien — pharmacien@sgp-pharma.tg", False),
         ("Caissier — caissier@sgp-pharma.tg", False),
         ("Magasinier — magasinier@sgp-pharma.tg", False),
         ("Sécurité", True),
         ("Mots de passe chiffrés (bcrypt)", False),
         ("Verrouillage après 5 tentatives (15 min)", False),
         ("Tokens JWT 8h + refresh 7 jours", False),
         ("Switch FR / EN en bas de la barre latérale", False),
     ],
     "01_dashboard.png"),

    ("MODULE 2 — DASHBOARD",
     "Tableau de bord",
     "Vue d'ensemble en temps réel : ventes, stock, alertes critiques.",
     [
         ("KPIs principaux (cartes vertes)", True),
         ("CA du jour & du mois (FCFA)", False),
         ("Valeur totale du stock + nombre d'unités", False),
         ("Nombre d'ordonnances servies aujourd'hui", False),
         ("Alertes critiques (rouge / orange)", True),
         ("Lots expirés (rouge) — à retirer immédiatement", False),
         ("Expirations < 30 jours (orange) — à écouler", False),
         ("Produits en stock faible (sous le seuil)", False),
         ("Visualisation", True),
         ("Graphique ventes des 7 derniers jours", False),
         ("Liste alertes péremption détaillée", False),
         ("Tableau produits sous le seuil d'alerte", False),
     ],
     "01_dashboard.png"),

    ("MODULE 3 — PRODUITS",
     "Produits — Catalogue CAMEG",
     "~150 médicaments essentiels couvrant tous les besoins du marché togolais.",
     [
         ("Données par produit", True),
         ("Code-barres, nom commercial, DCI", False),
         ("Forme pharmaceutique (cp, sirop, gélule…)", False),
         ("Catégorie (Antibiotiques, Antipaludiques…)", False),
         ("Prix de vente, seuil d'alerte stock", False),
         ("Indicateur ordonnance requise (Oui/Non)", False),
         ("Recherche intelligente", True),
         ("Recherche par nom, DCI ou code-barres", False),
         ("Catégories couvertes", True),
         ("Antalgiques, Antibiotiques, Antipaludiques", False),
         ("Cardiologie, Diabète, Gastro, Respiratoire", False),
         ("Dermato, Gynéco, Pédiatrie, Vitamines, Matériel", False),
         ("VIH/TB, Antiparasitaires, Anesthésiques, Solutions", False),
     ],
     "02_produits.png"),

    ("MODULE 4 — STOCK & LOTS",
     "Stock & Lots — Traçabilité FEFO",
     "Vue par lot avec date de péremption et badges colorés.",
     [
         ("Badges FEFO automatiques", True),
         ("OK (vert) — péremption > 90 jours", False),
         ("Bientôt (orange) — entre 0 et 90 jours", False),
         ("Expiré (rouge) — péremption dépassée", False),
         ("Informations par lot", True),
         ("Numéro de lot fabricant", False),
         ("Date de péremption + jours restants", False),
         ("Quantité actuelle / prix d'achat unitaire", False),
         ("Statut : Actif / Bloqué / Épuisé", False),
         ("Actions", True),
         ("Bloquer un lot (qualité douteuse)", False),
         ("Filtrage par produit ou texte libre", False),
         ("Tri automatique par péremption croissante", False),
     ],
     "03_stock_lots.png"),

    ("MODULE 5 — RÉCEPTION",
     "Réception fournisseur",
     "Entrée de stock multi-lignes avec validation stricte de péremption.",
     [
         ("Saisie obligatoire par ligne", True),
         ("Produit (sélection dans catalogue)", False),
         ("Numéro de lot fabricant", False),
         ("Date de péremption (refus si < aujourd'hui)", False),
         ("Prix d'achat unitaire", False),
         ("Quantité reçue", False),
         ("Règles métier", True),
         ("Blocage immédiat si lot déjà expiré", False),
         ("Création automatique mouvement ENTREE_ACHAT", False),
         ("Lien optionnel vers un bon de commande", False),
         ("Productivité", True),
         ("Réception multi-lignes en une seule transaction", False),
     ],
     "04_reception.png"),

    ("MODULE 6 — COMMANDES",
     "Bons de commande fournisseur",
     "Gestion des commandes avec génération PDF A4 imprimable.",
     [
         ("Cycle de commande", True),
         ("Brouillon → Validée → Reçue → Clôturée", False),
         ("Modification de statut depuis la liste", False),
         ("Génération PDF A4", True),
         ("Téléchargement instantané (icône télécharger)", False),
         ("En-tête vert OPTINET + nom pharmacie", False),
         ("Coordonnées fournisseur + date + statut", False),
         ("Tableau articles avec total HT en FCFA", False),
         ("Zone signature pharmacie + fournisseur", False),
         ("Données calculées", True),
         ("Total automatique (Σ qté × prix)", False),
         ("Audit log à chaque changement de statut", False),
     ],
     "05_commandes.png"),

    ("MODULE 7 — FOURNISSEURS",
     "Fournisseurs",
     "Annuaire des grossistes pharmaceutiques.",
     [
         ("Pré-chargés (démo)", True),
         ("CAMEG Togo (centrale d'achat nationale)", False),
         ("UBIPHARM Togo, COPHARMA, Laborex Togo", False),
         ("Champs gérés", True),
         ("Raison sociale, contact, email", False),
         ("Téléphone, adresse complète", False),
         ("Actions", True),
         ("Création / modification / suppression", False),
         ("Liaison aux lots (traçabilité origine)", False),
         ("Liaison aux bons de commande", False),
     ],
     "06_fournisseurs.png"),

    ("MODULE 8 — CAISSE POS",
     "Caisse (POS) — Vente avec FEFO automatique",
     "Interface optimisée pour le caissier — sélection de lot intelligente.",
     [
         ("Recherche & ajout panier", True),
         ("Recherche par nom / DCI / code-barres", False),
         ("Clic = ajout au panier (+1)", False),
         ("Indicateur ordonnance + stock disponible", False),
         ("Algorithme FEFO", True),
         ("Sélection auto du lot le plus proche péremption", False),
         ("Lots expirés exclus systématiquement", False),
         ("Décrémentation atomique (race-safe)", False),
         ("Ordonnance (si requise)", True),
         ("Référence texte + photo (upload base64)", False),
         ("Prévisualisation miniature avant validation", False),
         ("Ticket de caisse", True),
         ("PDF thermique 80mm téléchargeable", False),
         ("Modes : Espèces / Carte / Mobile money", False),
     ],
     "07_caisse_pos.png"),

    ("MODULE 9 — PERTES",
     "Pertes & Casses",
     "Justification fiscale et comptable des sorties non-vendeuses.",
     [
         ("3 motifs disponibles", True),
         ("Péremption — auto-détectable nightly", False),
         ("Casse — produit endommagé physiquement", False),
         ("Vol — disparition non justifiée", False),
         ("Effets sur le système", True),
         ("Décrémentation immédiate du lot concerné", False),
         ("Création de mouvement PERTE_* (immuable)", False),
         ("Audit log avec utilisateur + horodatage", False),
         ("Cas d'usage", True),
         ("Justification fiscale auprès du CAMEG", False),
         ("Suivi des pertes par motif sur la période", False),
         ("Identification produits à risque (vol récurrent)", False),
     ],
     "08_pertes.png"),

    ("MODULE 10 — RAPPORTS",
     "Rapports & Analyses",
     "3 rapports clés avec export CSV.",
     [
         ("Onglet Historique des ventes", True),
         ("Liste détaillée période (date, client, total, mode)", False),
         ("Export CSV pour comptabilité", False),
         ("Onglet Top produits", True),
         ("Classement quantité vendue + CA", False),
         ("Top 10 par défaut (configurable)", False),
         ("Onglet Marge brute", True),
         ("Calcul marge = CA − Coût (par produit)", False),
         ("Marge en FCFA + en pourcentage", False),
         ("Optimisation MongoDB aggregation $lookup", False),
     ],
     "09_rapports.png"),

    ("MODULE 11 — UTILISATEURS",
     "Gestion des utilisateurs",
     "Création de comptes + réinitialisation mots de passe sécurisée.",
     [
         ("Création utilisateur (Admin)", True),
         ("Email unique + nom complet", False),
         ("Mot de passe initial (≥ 6 caractères)", False),
         ("Rôle : Admin / Pharmacien / Caissier / Magasinier", False),
         ("Pharmacie d'affectation (Super Admin uniquement)", False),
         ("Réinitialisation mot de passe (icône clé)", True),
         ("Génération automatique pass temporaire fort", False),
         ("Affichage UNE SEULE FOIS dans une modale", False),
         ("Bouton de copie vers presse-papier", False),
         ("Activation / désactivation", True),
         ("Bouton Actif / Inactif (toggle)", False),
         ("Suppression avec confirmation", False),
     ],
     "10_utilisateurs.png"),

    ("MODULE 12 — AUDIT LOG",
     "Journal d'audit",
     "Traçabilité légale de toutes les actions critiques.",
     [
         ("Actions journalisées", True),
         ("Connexions / déconnexions", False),
         ("CRUD produits / suppliers / utilisateurs", False),
         ("Réceptions de stock", False),
         ("Ventes (montant, panier)", False),
         ("Pertes (motif, quantité)", False),
         ("Resets de mots de passe", False),
         ("Données capturées", True),
         ("Utilisateur (email + ID), horodatage UTC", False),
         ("Adresse IP, action, entité concernée, détails", False),
         ("Conformité", True),
         ("Conservation immuable (insert-only)", False),
     ],
     "11_audit_log.png"),

    ("MODULE 13 — MULTI-PHARMACIES",
     "Multi-pharmacies — Mode SaaS",
     "Architecture multi-tenant : OPTINET peut servir plusieurs officines.",
     [
         ("Rôle Super Admin (OPTINET)", True),
         ("Voit toutes les pharmacies du système", False),
         ("Crée / modifie / supprime des pharmacies", False),
         ("Crée des admins par pharmacie", False),
         ("Rôle Admin Pharmacie", True),
         ("Vue limitée à sa propre pharmacie", False),
         ("Affichage du nom en bandeau vert (sidebar)", False),
         ("Isolation stricte des données", True),
         ("Produits, stocks, ventes, utilisateurs cloisonnés", False),
         ("Cross-tenant access → erreur 404", False),
         ("Données par pharmacie", True),
         ("Nom, adresse, téléphone, email", False),
         ("N° d'agrément officiel, devise (FCFA)", False),
     ],
     "12_pharmacies.png"),

    ("VUE SUPER ADMIN",
     "Vue Super Admin OPTINET",
     "Tableau de bord agrégé sur l'ensemble du parc de pharmacies.",
     [
         ("Bandeau jaune Mode Super Admin", True),
         ("Visible dans la barre latérale gauche", False),
         ("Indique « Vue toutes pharmacies »", False),
         ("Données agrégées", True),
         ("CA cumulé toutes pharmacies confondues", False),
         ("Valeur stock totale du réseau", False),
         ("Alertes péremption consolidées", False),
         ("Cas d'usage", True),
         ("OPTINET supervise un réseau (Lomé, Kara…)", False),
         ("Onboarding rapide d'une nouvelle officine", False),
     ],
     "13_super_admin_dashboard.png"),
]


def add_filled_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_cover(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # Green background
    add_filled_rect(slide, 0, 0, prs.slide_width, prs.slide_height, PRIMARY)
    # Decorative cross
    add_filled_rect(slide, prs.slide_width - Inches(3.5), Inches(1), Inches(2.5), Inches(0.7), PRIMARY_DARK)
    add_filled_rect(slide, prs.slide_width - Inches(2.6), Inches(0.2), Inches(0.7), Inches(2.5), PRIMARY_DARK)

    add_text(slide, Inches(1), Inches(2.5), Inches(11.33), Inches(1.5),
             "SGP-Pharma", font_size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")
    add_text(slide, Inches(1), Inches(3.8), Inches(11.33), Inches(0.7),
             "Système de Gestion Intégrée de Pharmacie", font_size=24, color=RGBColor(0xA7, 0xF3, 0xD0), align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(5), Inches(11.33), Inches(0.5),
             "Manuel utilisateur — Présentation des modules", font_size=16, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(5.7), Inches(11.33), Inches(0.4),
             "OPTINET SARLU — Quartier Agoè Cacavéli, Lomé, Togo", font_size=13, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(6.1), Inches(11.33), Inches(0.4),
             "+228 90 74 84 65 · optinetsarl@gmail.com · www.optinet.tg", font_size=11, color=RGBColor(0xA7, 0xF3, 0xD0), align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(6.55), Inches(11.33), Inches(0.4),
             "Version 1.1 — Mai 2026", font_size=11, color=WHITE, align=PP_ALIGN.CENTER)


def add_module_slide(prs, badge, title, subtitle, bullets, screenshot):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Top green bar
    add_filled_rect(slide, 0, 0, prs.slide_width, Inches(0.4), PRIMARY)
    add_text(slide, Inches(0.4), Inches(0.08), Inches(7), Inches(0.3),
             "SGP-PHARMA · OPTINET SARLU", font_size=10, bold=True, color=WHITE)
    add_text(slide, Inches(prs.slide_width.inches - 4), Inches(0.08), Inches(3.5), Inches(0.3),
             "Manuel utilisateur · v1.1", font_size=10, color=WHITE, align=PP_ALIGN.RIGHT)

    # Badge
    add_text(slide, Inches(0.5), Inches(0.6), Inches(6), Inches(0.3),
             badge, font_size=10, bold=True, color=PRIMARY)
    # Title
    add_text(slide, Inches(0.5), Inches(0.9), Inches(6), Inches(0.7),
             title, font_size=28, bold=True, color=PRIMARY)
    # Subtitle
    add_text(slide, Inches(0.5), Inches(1.7), Inches(6), Inches(0.7),
             subtitle, font_size=12, color=MUTED)

    # Bullets
    bul_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(5.5), Inches(4.8))
    bul_tf = bul_box.text_frame
    bul_tf.word_wrap = True
    bul_tf.margin_left = bul_tf.margin_right = bul_tf.margin_top = bul_tf.margin_bottom = 0
    first = True
    for text, is_section in bullets:
        if first:
            p = bul_tf.paragraphs[0]
            first = False
        else:
            p = bul_tf.add_paragraph()
        run = p.add_run()
        if is_section:
            run.text = text
            run.font.bold = True
            run.font.size = Pt(13)
            run.font.color.rgb = TEXT
            p.space_before = Pt(6)
            p.space_after = Pt(2)
        else:
            run.text = "• " + text
            run.font.size = Pt(11)
            run.font.color.rgb = TEXT
            p.space_after = Pt(2)
        run.font.name = "Calibri"

    # Screenshot on right
    img_path = SCREENSHOTS / screenshot
    if img_path.exists():
        # 16:9 box on right side: width ~6.8 inches, height ~4.3 inches
        slide.shapes.add_picture(str(img_path), Inches(6.3), Inches(2.3), width=Inches(6.7))

    # Footer
    add_text(slide, Inches(0.5), Inches(prs.slide_height.inches - 0.4), Inches(7), Inches(0.3),
             "© 2026 OPTINET SARLU · Lomé, Togo", font_size=9, color=MUTED)


def add_toc(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_filled_rect(slide, 0, 0, prs.slide_width, Inches(0.4), PRIMARY)
    add_text(slide, Inches(0.4), Inches(0.08), Inches(7), Inches(0.3),
             "SGP-PHARMA · OPTINET SARLU", font_size=10, bold=True, color=WHITE)
    add_text(slide, Inches(0.5), Inches(0.6), Inches(6), Inches(0.3),
             "SOMMAIRE", font_size=10, bold=True, color=PRIMARY)
    add_text(slide, Inches(0.5), Inches(0.9), Inches(11), Inches(0.7),
             "Table des matières", font_size=28, bold=True, color=PRIMARY)
    add_text(slide, Inches(0.5), Inches(1.7), Inches(11), Inches(0.5),
             "13 modules + 1 vue Super Admin", font_size=12, color=MUTED)

    items = [
        ("1. Connexion & Rôles", "Sécurité JWT, 5 niveaux"),
        ("2. Tableau de bord", "KPIs + alertes péremption"),
        ("3. Produits", "Catalogue ~150 médicaments CAMEG"),
        ("4. Stock & Lots", "Gestion par lot avec FEFO"),
        ("5. Réception", "Entrée stock + validation"),
        ("6. Commandes fournisseur", "Bons + PDF A4"),
        ("7. Fournisseurs", "Annuaire CAMEG / Ubipharm"),
        ("8. Caisse (POS)", "Vente FEFO + ordonnance + ticket"),
        ("9. Pertes", "Péremption / Casse / Vol"),
        ("10. Rapports", "CA, top, marges, CSV"),
        ("11. Utilisateurs", "RBAC + reset mot de passe"),
        ("12. Journal d'audit", "Traçabilité légale"),
        ("13. Multi-pharmacies", "Mode SaaS Super Admin"),
        ("14. Vue Super Admin", "Tableau agrégé OPTINET"),
    ]
    y = Inches(2.5)
    col_w = Inches(6)
    for i, (n, d) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * col_w
        yy = y + Inches(0.5) * row
        add_text(slide, x, yy, Inches(3), Inches(0.4), n, font_size=12, bold=True, color=TEXT)
        add_text(slide, x + Inches(3), yy, Inches(3), Inches(0.4), d, font_size=10, color=MUTED)


def add_thanks(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_filled_rect(slide, 0, 0, prs.slide_width, prs.slide_height, PRIMARY)
    add_text(slide, Inches(1), Inches(1.5), Inches(11.33), Inches(1.5),
             "MERCI", font_size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(3), Inches(11.33), Inches(0.6),
             "Questions ?", font_size=24, color=RGBColor(0xA7, 0xF3, 0xD0), align=PP_ALIGN.CENTER)
    # Contact info
    add_text(slide, Inches(1), Inches(4), Inches(11.33), Inches(0.5),
             "OPTINET SARLU", font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(4.5), Inches(11.33), Inches(0.4),
             "Solutions Réseaux & Télécommunications", font_size=12, color=RGBColor(0xA7, 0xF3, 0xD0), align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(5.1), Inches(11.33), Inches(0.4),
             "Quartier Agoè Cacavéli, derrière la CEET — Lomé, Togo", font_size=12, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(5.5), Inches(11.33), Inches(0.4),
             "+228 90 74 84 65   ·   +228 99 05 84 71", font_size=12, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(5.9), Inches(11.33), Inches(0.4),
             "optinetsarl@gmail.com   ·   www.optinet.tg", font_size=12, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(6.4), Inches(11.33), Inches(0.4),
             "Directeur Général : NABINE Tassounti", font_size=11, color=RGBColor(0xA7, 0xF3, 0xD0), align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(6.85), Inches(11.33), Inches(0.4),
             "RCCM : TG-LFW-01-2026-B13-00831   ·   NIF : 1002114979", font_size=10, color=RGBColor(0xA7, 0xF3, 0xD0), align=PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_cover(prs)
    add_toc(prs)
    for badge, title, subtitle, bullets, ss in SLIDES:
        add_module_slide(prs, badge, title, subtitle, bullets, ss)
    add_thanks(prs)

    prs.save(str(OUT))
    size_kb = OUT.stat().st_size // 1024
    print(f"PowerPoint generated: {OUT} ({size_kb} KB, {len(SLIDES) + 3} slides)")


if __name__ == "__main__":
    build()
